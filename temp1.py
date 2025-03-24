import speech_recognition as sr
import pyttsx3
import os
import subprocess
import requests
import pyjokes
import smtplib
import spotipy
from spotipy.oauth2 import SpotifyClientCredentials
import nltk
from googletrans import Translator
import pandas as pd
import cv2
import pyautogui
import docx
import openpyxl
from pptx import Presentation
from PIL import ImageGrab
import hashlib
import datetime
import time

# Initialize speech recognition and synthesis
r = sr.Recognizer()
engine = pyttsx3.init()
engine.setProperty('rate', 150)  # Speed of speech
engine.setProperty('volume', 0.9)  # Volume level 0-1

USER_DATA_FILE = "user_data.csv"

# Function to load user data
def load_user_data():
    if os.path.exists(USER_DATA_FILE):
        return pd.read_csv(USER_DATA_FILE)
    else:
        return pd.DataFrame(columns=["username", "password"])

# Function to create a new user
def create_new_user():
    username = input("Enter new username: ")
    password = input("Enter new password: ")
    hashed_password = hashlib.sha256(password.encode()).hexdigest()
    new_user = pd.DataFrame([[username, hashed_password]], columns=["username", "password"])
    
    user_data = load_user_data()
    user_data = pd.concat([user_data, new_user], ignore_index=True)
    user_data.to_csv(USER_DATA_FILE, index=False)
    print("New user created successfully!")

# Check if user exists and authenticate
def authenticate_user():
    user_data = load_user_data()
    username = input("Enter your username: ")
    password = input("Enter your password: ")
    hashed_password = hashlib.sha256(password.encode()).hexdigest()
    
    if ((user_data["username"] == username) & (user_data["password"] == hashed_password)).any():
        print("Login successful!")
        return True
    else:
        print("Invalid username or password.")
        return False

# Function to open apps
def open_app(app_name):
    try:
        subprocess.Popen(app_name)
        print(f"Opened {app_name}")
    except Exception as e:
        print(f"Error opening {app_name}: {e}")

# Function to tell a joke
def tell_joke():
    joke = pyjokes.get_joke()
    print(joke)
    engine.say(joke)
    engine.runAndWait()

# Function to get IP address
def get_ip_address():
    ip_address = requests.get('https://api64.ipify.org').text
    print(f"Your IP address is: {ip_address}")
    return ip_address

# Function to get weather forecast
def get_weather_forecast():
    api_key = "69bf0a590576448ed0bfd804ac2b2694"  # replace with your OpenWeatherMap API key
    city = "New York"  # specify your city
    response = requests.get(f"http://api.openweathermap.org/data/2.5/weather?q={city}&appid={api_key}")
    weather_data = response.json()
    if weather_data.get("weather"):
        weather = weather_data["weather"][0]["description"]
        temperature = weather_data["main"]["temp"] - 273.15
        print(f"The weather in {city} is {weather} with a temperature of {temperature:.2f}°C")
        return weather, temperature
    else:
        print("Could not retrieve weather data.")
        return None

# Function to search on browser
def search_on_browser(query):
    url = f"https://www.google.com/search?q={query}"
    subprocess.Popen(['start', 'chrome', url], shell=True)

# Function to send message on WhatsApp (Placeholder)
def send_message_on_whatsapp(recipient, message):
    print(f"Sending message '{message}' to {recipient} on WhatsApp...")

# Function to install an app
def install_app(app_name):
    subprocess.run(f"winget install {app_name}", shell=True)
    print(f"{app_name} installed successfully.")

# Function to uninstall an app
def uninstall_app(app_name):
    subprocess.run(f"winget uninstall {app_name}", shell=True)
    print(f"{app_name} uninstalled successfully.")

# Function to download a file
def download_file(url):
    response = requests.get(url)
    filename = url.split("/")[-1]
    with open(filename, "wb") as file:
        file.write(response.content)
    print(f"{filename} downloaded successfully.")

# Function to record voice tasks
def record_voice_tasks():
    tasks = []
    with sr.Microphone() as source:
        print("Recording tasks, say 'stop' to finish...")
        while True:
            audio = r.listen(source)
            command = r.recognize_google(audio, language="en-in")
            if command.lower() == "stop":
                break
            tasks.append(command)
    tasks_df = pd.DataFrame(tasks, columns=["Task"])
    tasks_df.to_csv("voice_tasks.csv", index=False)
    print("Tasks recorded and saved to voice_tasks.csv")

# Face recognition placeholder
def face_recognition_auth():
    print("Authenticating with face recognition...")

# User authentication
def username_password_auth():
    if not authenticate_user():
        print("Authentication failed. Exiting.")
        return False
    return True

# Function to create and edit files
def create_and_edit_files(file_type, file_content):
    if file_type == "word":
        doc = docx.Document()
        doc.add_paragraph(file_content)
        doc.save("document.docx")
        print("Word document created.")
    elif file_type == "excel":
        workbook = openpyxl.Workbook()
        sheet = workbook.active
        sheet["A1"] = file_content
        workbook.save("spreadsheet.xlsx")
        print("Excel spreadsheet created.")
    elif file_type == "powerpoint":
        ppt = Presentation()
        slide = ppt.slides.add_slide(ppt.slide_layouts[0])
        title = slide.shapes.title
        title.text = file_content
        ppt.save("presentation.pptx")
        print("PowerPoint presentation created.")

# Function to send an email
def send_email(recipient, subject, body):
    try:
        server = smtplib.SMTP("smtp.gmail.com", 587)
        server.starttls()
        server.login("your_email@gmail.com", "your_password")
        message = f"Subject: {subject}\n\n{body}"
        server.sendmail("your_email@gmail.com", recipient, message)
        server.quit()
        print("Email sent successfully.")
    except Exception as e:
        print(f"Error sending email: {e}")

# Function to update drivers
def update_drivers():
    print("Updating drivers...")
    os.system("msdt.exe /id DeviceDiagnostic")

# Function to play music from Spotify
def play_music_from_spotify():
    sp = spotipy.Spotify(auth_manager=SpotifyClientCredentials(client_id="your_spotify_client_id", client_secret="your_spotify_client_secret"))
    results = sp.current_user_playing_track()
    if results:
        track = results["item"]["name"]
        print(f"Playing {track}")
    else:
        print("No track is currently playing.")

# Function to take a screenshot
def take_screenshot():
    screenshot = ImageGrab.grab()
    screenshot.save("screenshot.png")
    print("Screenshot saved as screenshot.png")

# Function to record a video
def record_video():
    cap = cv2.VideoCapture(0)
    out = cv2.VideoWriter("video.avi", cv2.VideoWriter_fourcc(*'XVID'), 20.0, (640, 480))
    while cap.isOpened():
        ret, frame = cap.read()
        if ret:
            out.write(frame)
            cv2.imshow("Recording...", frame)
            if cv2.waitKey(1) & 0xFF == ord("q"):
                break
        else:
            break
    cap.release()
    out.release()
    cv2.destroyAllWindows()
    print("Video recording saved as video.avi")

# Function to translate language
def translate_language(text, target_language):
    translator = Translator()
    translation = translator.translate(text, dest=target_language)
    print(f"Translated text: {translation.text}")
    return translation.text

# Function to summarize text
def summarize_text(text):
    sentences = nltk.sent_tokenize(text)
    summary = " ".join(sentences[:2])
    print(f"Summary: {summary}")
    return summary

# Main loop
def main():
    if not os.path.exists(USER_DATA_FILE) or input("Are you a new user? (yes/no): ").strip().lower() == "yes":
        create_new_user()
    else:
        if not authenticate_user():
            print("Authentication failed. Exiting.")
            return

    while True:
        with sr.Microphone() as source:
            print("Listening for commands...")
            audio = r.listen(source)
            try:
                command = r.recognize_google(audio, language="en-in")
                print(f"You said: {command}")
                if "joke" in command:
                    tell_joke()
                elif "IP address" in command:
                    get_ip_address()
                elif "weather" in command:
                    get_weather_forecast()
                elif "search" in command:
                    search_on_browser(command.replace("search", "").strip())
                elif "open" in command:
                    open_app(command.replace("open", "").strip())
                elif "screenshot" in command:
                    take_screenshot()
                elif "record video" in command:
                    record_video()
                elif "email" in command:
                    recipient = input("Enter the recipient email address: ")
                    subject = input("Enter the subject: ")
                    body = input("Enter the email body: ")
                    send_email(recipient, subject, body)
                elif "Spotify" in command:
                    play_music_from_spotify()
                elif "install" in command:
                    app_name = command.replace("install", "").strip()
                    install_app(app_name)
                elif "uninstall" in command:
                    app_name = command.replace("uninstall", "").strip()
                    uninstall_app(app_name)
                elif "WhatsApp" in command:
                    recipient = input("Enter the recipient's name: ")
                    message = input("Enter your message: ")
                    send_message_on_whatsapp(recipient, message)
                elif "translate" in command:
                    text = input("Enter the text to translate: ")
                    target_language = input("Enter the target language code (e.g., 'es' for Spanish): ")
                    translate_language(text, target_language)
                elif "summarize" in command:
                    text = input("Enter the text to summarize: ")
                    summarize_text(text)
                elif "download" in command:
                    url = input("Enter the file URL to download: ")
                    download_file(url)
                elif "record tasks" in command:
                    record_voice_tasks()
                elif "create file" in command:
                    file_type = input("Enter file type (word/excel/powerpoint): ")
                    content = input("Enter content for the file: ")
                    create_and_edit_files(file_type, content)
                elif "drivers" in command:
                    update_drivers()
                elif "exit" in command:
                    print("Goodbye!")
                    break
                else:
                    print("Command not recognized. Please try again.")
            except sr.UnknownValueError:
                print("Sorry, I didn't understand that. Please speak clearly.")
            except Exception as e:
                print(f"Error: {e}")

if "_name_" == "_main_":
    main()